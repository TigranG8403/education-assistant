from kivy.app import App
from kivy.uix.label import Label
from kivy.uix.button import Button
from kivy.uix.recycleview import RecycleView
from kivy.uix.popup import Popup
from kivy.uix.textinput import TextInput
from kivy.uix.screenmanager import ScreenManager, Screen
from kivy.graphics import Color, Rectangle, Line
from kivy.config import Config
from kivy.uix.image import Image
from kivy.lang import Builder
from kivy.clock import Clock
from pydantic import BaseModel
import requests
import speech_recognition as sr
import pyaudio
from threading import Thread
import wave
from pydub import AudioSegment
from moviepy.editor import VideoFileClip
import os
import fitz
from docx import Document
from pptx import Presentation
from functools import partial


sm = ScreenManager()

Config.set('postproc', 'file', '')
Config.write()

Builder.load_string('''

<RV>:
    viewclass: 'Button'
    RecycleBoxLayout:
        default_size: None, dp(56)
        default_size_hint: 1, None
        size_hint_y: None
        height: self.minimum_height + 0.2*self.minimum_height
        orientation: 'vertical'

<FilesView>:
    cols:1
    id: FilesView
    
    FileChooserListView:
        id: filechooser
        
        filters: ['*.mp4', '*.avi']
        on_selection:FilesView.selected(filechooser.selection)

<FilesView1>:
    cols:1
    id: FilesView
    
    FileChooserListView:
        id: filechooser
        filters: ['*.wav']
        on_selection:FilesView.selected(filechooser.selection)

<FilesView2>:
    cols:1
    id: FilesView
    
    FileChooserListView:
        id: filechooser
        filters: ['*.txt', '*.doc', '*.docx', '*.pdf', '*.pptx']
        on_selection:FilesView.selected(filechooser.selection)

''')



class RV(RecycleView):
    def __init__(self, **kwargs):
        super(RV, self).__init__(**kwargs)
        user_folder = Config.get('postproc', 'email')
        files=[]
        response = requests.get(f"http://176.56.54.46:1234/files/{user_folder}")
        if response.status_code == 200:
            data1 = response.json()
            for filename in data1["files"]:
                files.append(filename)
            
            self.data = [{'text': str(filename), 'on_release': partial(self.butt, filename)} for filename in files]
            
    def butt(self, filename):
        if filename:
                
                Config.set('postproc', 'file', filename)
                Config.write()
                sm.current = 'result'
                

class FilesView(Popup):
    def selected(self, filename):
            if filename:
                file = filename[0]
                self.dismiss()
                extract_audio_from_video(file, 'extracted_audio.wav')

class FilesView1(Popup):
    def selected(self, filename1):
        if filename1:
            audio_file = filename1[0]
            # Закрыть popup
            self.dismiss()
            # Запустить функцию для обработки файла
            split_recog(audio_file)

class FilesView2(Popup):
    def selected(self, filename2):
        if filename2:
            # Закрыть popup
            self.dismiss()
            # Запустить функцию для обработки файла
            extract_text(filename2[0])
                
def extract_text(filename2):
  _, extension = os.path.splitext(filename2)

  if extension == ".pdf":
    return extract_text_from_pdf(filename2)
  elif extension == ".docx":
    return extract_text_from_docx(filename2)
  elif extension == ".pptx":
    return extract_text_from_pptx(filename2)
  elif extension == ".txt":
    return extract_text_from_txt(filename2)
      
def extract_text_from_txt(filename2):
    file = open(filename2, "r")
    text = file.read()
    file.close()
    Config.set('postproc', 'text', text)
    Config.write()
    sm.current = 'result'
    

def extract_text_from_pdf(filename2):
  with fitz.open(filename2) as doc:
    text = ""
    for page in doc:
      text += page.get_text()
    Config.set('postproc', 'text', text)
    Config.write()
    sm.current = 'result'
    

def extract_text_from_docx(filename2):
  doc = Document(filename2)
  text = ""
  for paragraph in doc.paragraphs:
    text += paragraph.text + "\n"
  Config.set('postproc', 'text', text)
  Config.write()
  sm.current = 'result'
  

def extract_text_from_pptx(filename2):
  prs = Presentation(filename2)
  text = ""
  for slide in prs.slides:
    for shape in slide.shapes:
      if hasattr(shape, "text"):
        text += shape.text + "\n"
  Config.set('postproc', 'text', text)
  Config.write()
  sm.current = 'result'
  

class RegScreen(Screen):
    def __init__(self, **kwargs):
        super().__init__(**kwargs)

        self.name = 'reg'

        with self.canvas:
            # Рисуем прямоугольник для градиента 1
            Color(1, 0.84, 0)  # Золотой
            self.rect1 = Rectangle(pos=(0, 0), size=(self.width / 2, self.height))

            # Рисуем прямоугольник для градиента 2
            Color(1, 0.2, 0.2)  # Красный
            self.rect2 = Rectangle(pos=(self.width / 2, 0), size=(self.width / 2, self.height))

            # Рисуем линии, формирующие пересечение
            self.line1 = Line(points=(0, 0), width=self.width, close=True, cap='round')
            self.line2 = Line(points=(0, self.height), width=self.height, close=True, cap='round')

        # Привязываем размеры прямоугольников и линий к размеру виджета
        self.bind(size=self._update_background)
        # Текст "COPY-PASTE"
        self.add_widget(Label(text="COPY-PASTE",
                              font_size="24sp",
                              pos_hint={"center_x": 0.5, "center_y": 0.9},
                              halign="center"))

        # Текст "Регистрация"
        self.add_widget(Label(text="Регистрация",
                              font_size="36sp",
                              pos_hint={"center_x": 0.5, "center_y": 0.7},
                              halign="center"))

        # Поле ввода "Email"
        self.email_input = TextInput(hint_text="Email",
                                   multiline=False,
                                   pos_hint={"center_x": 0.5, "center_y": 0.6},
                                   size_hint=(0.8, None),
                                   height="40sp",
                                   halign="center")
        self.add_widget(self.email_input)

        # Поле ввода "Пароль"
        self.password_input = TextInput(hint_text="Пароль",
                                      multiline=False,
                                      pos_hint={"center_x": 0.5, "center_y": 0.5},
                                      size_hint=(0.8, None),
                                      height="40sp",
                                      halign="center",
                                      password=True)
        self.add_widget(self.password_input)

        # Кнопка "Войти" (btn1)
        self.btn1 = Button(text="Уже есть аккаунт? Войти",
                           font_size="18sp",
                           pos_hint={"center_x": 0.5, "center_y": 0.4},
                           size_hint=(0.265, 0.05),
                           background_color=(1, 1, 1, 0),  # Прозрачный фон
                           border=(0, 0, 0, 0),  # Удаляем границы
                           color=(0, 0, 1, 1),  # Цвет текста
                           halign="center")
        self.btn1.bind(on_press=self.on_btn1_press)
        self.add_widget(self.btn1)

        # Кнопка "Че угодно" (btn2)
        self.btn2 = Button(text="Подтвердить",
                           pos_hint={"center_x": 0.5, "center_y": 0.2},
                           size_hint=(0.265, None),
                           height="40sp")
        self.btn2.bind(on_press=self.on_btn2_press)
        self.add_widget(self.btn2)

    def on_btn1_press(self, instance):
        self.manager.current = 'auth'
        

    def on_btn2_press(self, instance):
        email = self.email_input.text
        password = self.password_input.text

        user = {"email": email, "password": password}
        response = requests.post("http://176.56.54.46:1234/register", json=user)

        if response.status_code == 200:
            self.manager.current = 'auth'
        else:
            popup = Popup(title='Ошибка регистрации')
            popup.content = Label(text=response.text)
            popup.size_hint = (None, None)
            popup.size = (400, 400)
            popup.open()


    def _update_background(self, instance, value):
        # Обновляем размеры прямоугольников и линий при изменении размера виджета
        self.rect1.size = value
        self.rect2.size = value
        self.line1.width = value[0]
        self.line2.width = value[1]
        self.line2.points = (0, value[1], value[0], value[1], value[0], 0)  # Обновляем точки второй линии

class AuthScreen(Screen):
    def __init__(self, **kwargs):
        super().__init__(**kwargs)
        
        self.name = 'auth'

        with self.canvas:
            # Рисуем прямоугольник для градиента 1
            Color(1, 0.84, 0)  # Золотой
            self.rect1 = Rectangle(pos=(0, 0), size=(self.width / 2, self.height))

            # Рисуем прямоугольник для градиента 2
            Color(1, 0.2, 0.2)  # Красный
            self.rect2 = Rectangle(pos=(self.width / 2, 0), size=(self.width / 2, self.height))

            # Рисуем линии, формирующие пересечение
            self.line1 = Line(points=(0, 0), width=self.width, close=True, cap='round')
            self.line2 = Line(points=(0, self.height), width=self.height, close=True, cap='round')

        # Привязываем размеры прямоугольников и линий к размеру виджета
        self.bind(size=self._update_background)
        # Текст "COPY-PASTE"
        self.add_widget(Label(text="COPY-PASTE",
                              font_size="24sp",
                              pos_hint={"center_x": 0.5, "center_y": 0.9},
                              halign="center"))

        # Текст "Регистрация"
        self.add_widget(Label(text="Авторизация",
                              font_size="36sp",
                              pos_hint={"center_x": 0.5, "center_y": 0.7},
                              halign="center"))

        # Поле ввода "Email"
        self.email_input = TextInput(hint_text="Email",
                                   multiline=False,
                                   pos_hint={"center_x": 0.5, "center_y": 0.6},
                                   size_hint=(0.8, None),
                                   height="40sp",
                                   halign="center")
        self.add_widget(self.email_input)

        # Поле ввода "Пароль"
        self.password_input = TextInput(hint_text="Пароль",
                                      multiline=False,
                                      pos_hint={"center_x": 0.5, "center_y": 0.5},
                                      size_hint=(0.8, None),
                                      height="40sp",
                                      halign="center",
                                      password=True)
        self.add_widget(self.password_input)

        # Кнопка "Войти" (btn1)
        self.btn1 = Button(text="Ещё нет аккаунта? Создать",
                           font_size="18sp",
                           pos_hint={"center_x": 0.5, "center_y": 0.4},
                           size_hint=(0.265, 0.05),
                           background_color=(1, 1, 1, 0),  # Прозрачный фон
                           border=(0, 0, 0, 0),  # Удаляем границы
                           color=(0, 0, 1, 1),  # Цвет текста
                           halign="center")
        self.btn1.bind(on_press=self.on_btn1_press)
        self.add_widget(self.btn1)

        # Кнопка "Че угодно" (btn2)
        self.btn2 = Button(text="Войти",
                           pos_hint={"center_x": 0.5, "center_y": 0.2},
                           size_hint=(0.265, None),
                           height="40sp")
        self.btn2.bind(on_press=self.on_btn2_press)
        self.add_widget(self.btn2)

    def on_btn1_press(self, instance):
        self.manager.current = 'reg'  # Вывод сообщения в консоль при нажатии на btn1

    def on_btn2_press(self, instance):
        email = self.email_input.text
        password = self.password_input.text
        
        user = User(email=email, password=password)

        response = requests.post("http://176.56.54.46:1234/auth", json=user.dict())

        if response.status_code == 200:
            Config.set('postproc', 'email', email)
            Config.write()
            Config.set('postproc', 'password', password)
            Config.write()
            self.manager.current = 'main'
        else:
            popup = Popup(title='Ошибка авторизации')
            popup.content = Label(text=response.text)
            popup.size_hint = (None, None)
            popup.size = (400, 400)
            popup.open()



    def _update_background(self, instance, value):
        self.rect1.size = value
        self.rect2.size = value
        self.line1.width = value[0]
        self.line2.width = value[1]
        self.line2.points = (0, value[1], value[0], value[1], value[0], 0)  # Обновляем точки второй линии

class MainScreen(Screen):
    def __init__(self, **kwargs):
        super().__init__(**kwargs)
        self.name = 'main'
        
        with self.canvas:
            Color(0.9, 0.2, 0.2, 1)
            self.rect_top = Rectangle(pos=(0, 0), size=(self.width, self.height / 2))
            Color(0.9, 0.6, 0.2, 1)
            self.rect_bottom = Rectangle(pos=(0, self.height / 2), size=(self.width, self.height / 2))
            Color(0, 0, 0, 1)
            Line(points=[(self.width * 0.05, self.height * 0.05), (self.width * 0.1, self.height * 0.05),
                         (self.width * 0.1, self.height * 0.15), (self.width * 0.05, self.height * 0.15)])
            Line(points=[(self.width * 0.1, self.height * 0.05), (self.width * 0.1, self.height * 0.15),
                         (self.width * 0.15, self.height * 0.15), (self.width * 0.15, self.height * 0.05),
                         (self.width * 0.2, self.height * 0.05), (self.width * 0.2, self.height * 0.15)])
            Line(points=[(self.width * 0.2, self.height * 0.05), (self.width * 0.2, self.height * 0.15),
                         (self.width * 0.25, self.height * 0.15), (self.width * 0.25, self.height * 0.05)])

        title = Label(text='Education Assistant', font_size=24, bold=True,
                      pos_hint={'center_x': 0.5, 'center_y': 0.9}, color=[0, 0, 0, 1])
        self.add_widget(title)

        settings = Image(source='settings.png',
                      pos_hint={'center_x': 0.1, 'center_y': 0.9},
                      size_hint=(0.1, 0.1), allow_stretch=True)
        self.add_widget(settings)

        self.settings_button = Button(pos_hint={'center_x': 0.1, 'center_y': 0.9},
                           size_hint=(0.1, 0.1),
                           background_color=(1, 1, 1, 0),  # Прозрачный фон
                           border=(0, 0, 0, 0),  # Удаляем границы
                           color=(0, 0, 1, 1) # Цвет текста
                          )
        self.settings_button.bind(on_press=self.on_settings_button_pressed)
        self.add_widget(self.settings_button)

        folder = Image(source='folder.png',
                      pos_hint={'center_x': 0.9, 'center_y': 0.9},
                      size_hint=(0.1, 0.1), allow_stretch=True)
        self.add_widget(folder)

        self.folder_button = Button(pos_hint={'center_x': 0.9, 'center_y': 0.9},
                           size_hint=(0.1, 0.1),
                           background_color=(1, 1, 1, 0),  # Прозрачный фон
                           border=(0, 0, 0, 0),  # Удаляем границы
                           color=(0, 0, 1, 1) # Цвет текста
                          )
        self.folder_button.bind(on_press=self.on_folder_button_pressed)
        self.add_widget(self.folder_button)


        input_label = Label(text='Введите текст', pos_hint={'center_x': 0.5, 'center_y': 0.7},
                            color=[0, 0, 0, 1]
                            )
        self.add_widget(input_label)

        self.input_field = TextInput(hint_text="Введите текст или ссылку на Ютуб",
                                   size_hint=(0.8, None),
                                   pos_hint={'center_x': 0.5, 'center_y': 0.7},
                                   readonly=False,
                                   halign='left', padding=[20, 20, 20, 20])
        self.add_widget(self.input_field)

        self.send_button = Button(
                      background_color=(1, 1, 1, 0),
                      border=(0, 0, 0, 0),
                      color=(0, 0, 1, 1),
                      pos_hint={'center_x': 0.9, 'center_y': 0.7},
                      size_hint=(0.15, None))
        self.send_button.bind(on_press=self.on_send_button_pressed)
        self.add_widget(self.send_button)

        self.send_image = Image(source='send.png',
                     pos_hint={'center_x': 0.9, 'center_y': 0.7},
                     size_hint=(0.15, None), allow_stretch=True
            )
        self.add_widget(self.send_image)

        

        or_label = Label(text='ИЛИ', pos_hint={'center_x': 0.5, 'center_y': 0.55}, # Изменили center_y
                          color=[0, 0, 0, 1], font_size=20)
        self.add_widget(or_label)

        # Создаем кнопки и labels сразу
        self.record_button = Button(background_color=(1, 1, 1, 0),
                      border=(0, 0, 0, 0),
                      color=(0, 0, 1, 1),
                              pos_hint={'center_x': 0.5, 'center_y': 0.4},
                              size_hint=(None, None))
        self.record_button.bind(on_press=self.start_recording)
        self.add_widget(self.record_button)

        self.record_image = Image(source='record.png',
                     pos_hint={'center_x': 0.5, 'center_y': 0.4},
                     size_hint=(None, None)
            )
        self.add_widget(self.record_image)

        record_label = Image(source='record_label.png',
                     pos_hint={'center_x': 0.5, 'center_y': 0.3},
                     size_hint=(None, None)
            )
        self.add_widget(record_label)
        

        self.video_button = Button(background_color=(1, 1, 1, 0),
                      border=(0, 0, 0, 0),
                      color=(0, 0, 1, 1),
                             pos_hint={'center_x': 0.2, 'center_y': 0.2},
                             size_hint=(None, None))
        self.add_widget(self.video_button)
        self.video_button.bind(on_press=self.video_button_pressed)

        video_image = Image(source='video.png',
                     pos_hint={'center_x': 0.2, 'center_y': 0.2},
                     size_hint=(None, None)
            )
        self.add_widget(video_image)

        video_label = Image(source='video_label.png',
                     pos_hint={'center_x': 0.2, 'center_y': 0.06},
                     size_hint=(None, None)
            )
        self.add_widget(video_label)

        self.docx_button = Button(background_color=(1, 1, 1, 0),
                      border=(0, 0, 0, 0),
                      color=(0, 0, 1, 1),
                              pos_hint={'center_x': 0.5, 'center_y': 0.2},
                              size_hint=(None, None))
        self.docx_button.bind(on_press=self.on_docx_button_pressed)
        self.add_widget(self.docx_button)


        docx_button_image = Image(source='docx.png',
                     pos_hint={'center_x': 0.5, 'center_y': 0.2},
                     size_hint=(None, None)
            )
        self.add_widget(docx_button_image)

        docx_label = Image(source='docx_label.png',
                     pos_hint={'center_x': 0.5, 'center_y': 0.06},
                     size_hint=(None, None)
            )
        self.add_widget(docx_label)


        self.audio_button = Button(
            background_color=(1, 1, 1, 0),
                      border=(0, 0, 0, 0),
                      color=(0, 0, 1, 1),
            pos_hint={'center_x': 0.8, 'center_y': 0.2},
            size_hint=(None, None)
        )
        self.add_widget(self.audio_button)
        self.audio_button.bind(on_press=self.audio_button_pressed)

        audio_image = Image(source='audio.png',
                     pos_hint={'center_x': 0.8, 'center_y': 0.2},
                     size_hint=(None, None)
            )
        self.add_widget(audio_image)

        audio_label = Image(source='audio_label.png',
                     pos_hint={'center_x': 0.8, 'center_y': 0.06},
                     size_hint=(None, None)
            )
        self.add_widget(audio_label)

        self.stop_record_button = Button(background_color=(1, 1, 1, 0),
                      border=(0, 0, 0, 0),
                      color=(0, 0, 1, 1),
                              pos_hint={'center_x': 0.5, 'center_y': 0.4},
                              size_hint=(None, None))
        self.stop_record_button.bind(on_press=self.stop_record_button_pressed)
        self.stop_record_button.pressed = False
        

        self.stop_record_image = Image(source='stop_record.png',
                     pos_hint={'center_x': 0.5, 'center_y': 0.4},
                     size_hint=(None, None)
            )

    def delete_widget(self):
        try:
            self.remove_widget(self.stop_record_image)
            self.remove_widget(self.stop_record_button)
            self.add_widget(self.record_image)
            self.add_widget(self.record_button)
        except: 
            pass

        

    def on_size(self, *args):
        self.rect_top.size = self.size
        self.rect_bottom.pos = (0, self.height / 2)
        self.rect_bottom.size = self.size

    def on_settings_button_pressed(self, instance):
        self.manager.current = 'settings'

    def on_folder_button_pressed(self, instance):
        
        self.manager.get_screen('archive').delete_widget()

        self.manager.current = 'archive'
        return

    def recording(self):
        
        FORMAT = pyaudio.paInt16  # Формат аудиоданных
        CHANNELS = 1  # Количество каналов
        RATE = 44100  # Частота дискретизации (Гц)
        CHUNK = 1024  # Размер блока данных
        RECORD_SECONDS = 100000000  # Длительность записи (секунды)
        
        p = pyaudio.PyAudio()
        
        stream = p.open(
        format=FORMAT,
        channels=CHANNELS,
        rate=RATE,
        input=True,
        frames_per_buffer=CHUNK
        )
        
        frames = []

        for j in range(0, int(RATE / CHUNK * RECORD_SECONDS)):
            if not self.stop_record_button.pressed:
                    data = stream.read(CHUNK)
                    frames.append(data)

            else:
                break

        print("Запись завершена.")

        stream.stop_stream()
        stream.close()

        wf = wave.open('audio_cache.wav', 'wb')
        wf.setnchannels(CHANNELS)
        wf.setsampwidth(p.get_sample_size(FORMAT))
        wf.setframerate(RATE)
        wf.writeframes(b''.join(frames))
        wf.close()
        thread1 = Thread(target=split_recog('audio_cache.wav'))
        thread1.start()

    def start_recording(self, instance):
        self.record = True
        self.recording = self.recording
        self.remove_widget(self.record_button)
        self.remove_widget(self.record_image)
        self.add_widget(self.stop_record_image)
        self.add_widget(self.stop_record_button)
        thread = Thread(target=self.recording)
        thread.start()
        
    def stop_record_button_pressed(self, instance):
        self.stop_record_button.pressed = True
        

    def video_button_pressed(self, instance):
        popup = FilesView(size_hint = (0.8, 0.8))
        popup.open()
        return
        

    def audio_button_pressed(self, instance):
        popup = FilesView1(size_hint = (0.8, 0.8))
        popup.open()
        return
        
    def on_docx_button_pressed(self, instance):
        popup = FilesView2(size_hint = (0.8, 0.8))
        popup.open()
        return

    def on_send_button_pressed(self, instance):
        prefixes = [
  "youtube.com","www.youtube.com",
  "https://www.youtube.com",
  "http://www.youtube.com",
  "youtu.be",
  "www.youtu.be",
  "https://www.youtu.be",
  "http://www.youtu.be",
  "https://youtu.be",
  "http://youtu.be",
  "short.youtube.com",
  "www.short.youtube.com",
  "https://www.short.youtube.com",
  "http://www.short.youtube.com",
  "music.youtube.com",
  "www.music.youtube.com",
  "https://www.music.youtube.com",
  "http://www.music.youtube.com",
  "studio.youtube.com",
  "www.studio.youtube.com",
  "https://www.studio.youtube.com",
  "http://www.studio.youtube.com",
  "gaming.youtube.com",
  "www.gaming.youtube.com",
  "https://www.gaming.youtube.com",
  "http://www.gaming.youtube.com",
  "tv.youtube.com",
  "www.tv.youtube.com",
  "https://www.tv.youtube.com",
  "http://www.tv.youtube.com",
  "kids.youtube.com",
  "www.kids.youtube.com",
  "https://www.kids.youtube.com",
  "http://www.kids.youtube.com",
  "youtube.com/creators",
  "www.youtube.com/creators",
  "https://www.youtube.com/creators",
  "http://www.youtube.com/creators",
  "youtube.com/activate",
  "www.youtube.com/activate",
  "https://www.youtube.com/activate",
  "http://www.youtube.com/activate",
  "youtube.co.uk",
  "www.youtube.co.uk",
  "https://www.youtube.co.uk",
  "http://www.youtube.co.uk",
  "youtube.de",
  "www.youtube.de",
  "https://www.youtube.de",
  "http://www.youtube.de",
  "youtube.fr",
  "www.youtube.fr",
  "https://www.youtube.fr",
  "http://www.youtube.fr",
  "youtube.es",
  "www.youtube.es",
  "https://www.youtube.es",
  "http://www.youtube.es",
  "youtube.it",
  "www.youtube.it",
  "https://www.youtube.it",
  "http://www.youtube.it",
  "youtube.ca",
  "www.youtube.ca",
  "https://www.youtube.ca",
  "http://www.youtube.ca",
  "youtube.com.br",
  "www.youtube.com.br",
  "https://www.youtube.com.br",
  "http://www.youtube.com.br",
  "youtube.com.au",
  "www.youtube.com.au",
  "https://www.youtube.com.au",
  "http://www.youtube.com.au",
  "youtube.com.mx",
  "www.youtube.com.mx",
  "https://www.youtube.com.mx",
  "http://www.youtube.com.mx",
  "youtube.com.tr",
  "www.youtube.com.tr",
  "https://www.youtube.com.tr",
  "http://www.youtube.com.tr",
  "youtube.co.jp",
  "www.youtube.co.jp",
  "https://www.youtube.co.jp",
  "http://www.youtube.co.jp",
  "youtube.com.hk",
  "www.youtube.com.hk",
  "https://www.youtube.com.hk",
  "http://www.youtube.com.hk",
  "youtube.com.tw",
  "www.youtube.com.tw",
  "https://www.youtube.com.tw",
  "http://www.youtube.com.tw",
  "youtube.com.sg",
  "www.youtube.com.sg",
  "https://www.youtube.com.sg",
  "http://www.youtube.com.sg",
  "youtube.com.my",
  "www.youtube.com.my",
  "https://www.youtube.com.my",
  "http://www.youtube.com.my",
  "youtube.com.ph",
  "www.youtube.com.ph",
  "https://www.youtube.com.ph",
  "http://www.youtube.com.ph",
  "youtube.com.vn",
  "www.youtube.com.vn",
  "https://www.youtube.com.vn",
  "http://www.youtube.com.vn",
  "youtube.co.in",
  "www.youtube.co.in",
  "https://www.youtube.co.in",
  "http://www.youtube.co.in",
  "youtube.co.kr",
  "www.youtube.co.kr",
  "https://www.youtube.co.kr",
  "http://www.youtube.co.kr",
  "youtube.ru",
  "www.youtube.ru",
  "https://www.youtube.ru",
  "http://www.youtube.ru",
  "youtube.com.ua",
  "www.youtube.com.ua",
  "https://www.youtube.com.ua",
  "http://www.youtube.com.ua",
  "youtube.pl",
  "www.youtube.pl",
  "https://www.youtube.pl",
  "http://www.youtube.pl",
  "youtube.cz",
  "www.youtube.cz",
  "https://www.youtube.cz",
  "http://www.youtube.cz",
  "youtube.sk",
  "www.youtube.sk",
  "https://www.youtube.sk",
  "http://www.youtube.sk",
  "youtube.ro",
  "www.youtube.ro",
  "https://www.youtube.ro",
  "http://www.youtube.ro",
  "youtube.hu",
  "www.youtube.hu",
  "https://www.youtube.hu",
  "http://www.youtube.hu",
  "youtube.gr",
  "www.youtube.gr",
  "https://www.youtube.gr",
  "http://www.youtube.gr",
  "youtube.pt",
  "www.youtube.pt",
  "https://www.youtube.pt",
  "http://www.youtube.pt",
  "youtube.se",
  "www.youtube.se",
  "https://www.youtube.se",
  "http://www.youtube.se",
  "youtube.fi",
  "www.youtube.fi",
  "https://www.youtube.fi",
  "http://www.youtube.fi",
  "youtube.no",
  "www.youtube.no",
  "https://www.youtube.no",
  "http://www.youtube.no",
  "youtube.dk",
  "www.youtube.dk",
  "https://www.youtube.dk",
  "http://www.youtube.dk",
  "youtube.be",
  "www.youtube.be",
  "https://www.youtube.be",
  "http://www.youtube.be",
  "youtube.nl",
  "www.youtube.nl",
  "https://www.youtube.nl",
  "http://www.youtube.nl",
  "youtube.ch",
  "www.youtube.ch",
  "https://www.youtube.ch",
  "http://www.youtube.ch",
  "youtube.at",
  "www.youtube.at",
  "https://www.youtube.at",
  "http://www.youtube.at",
  "youtube.com.ar",
  "www.youtube.com.ar",
  "https://www.youtube.com.ar",
  "http://www.youtube.com.ar",
  "youtube.com.co",
  "www.youtube.com.co",
  "https://www.youtube.com.co",
  "http://www.youtube.com.co",
  "youtube.com.pe",
  "www.youtube.com.pe",
  "https://www.youtube.com.pe",
  "http://www.youtube.com.pe",
  "youtube.com.ve",
  "www.youtube.com.ve",
  "https://www.youtube.com.ve",
  "http://www.youtube.com.ve",
  "youtube.com.ec",
  "www.youtube.com.ec",
  "https://www.youtube.com.ec",
  "http://www.youtube.com.ec",
  "youtube.com.cl",
  "www.youtube.com.cl",
  "https://www.youtube.com.cl",
  "http://www.youtube.com.cl",
  "youtube.com.uy",
  "www.youtube.com.uy",
  "https://www.youtube.com.uy",
  "http://www.youtube.com.uy",
  "youtube.com.py",
  "www.youtube.com.py",
  "https://www.youtube.com.py",
  "http://www.youtube.com.py",
  "youtube.com.sa",
  "www.youtube.com.sa",
  "https://www.youtube.com.sa",
  "http://www.youtube.com.sa",
  "youtube.com.ae",
  "www.youtube.com.ae",
  "https://www.youtube.com.ae",
  "http://www.youtube.com.ae",
  "youtube.com.eg",
  "www.youtube.com.eg",
  "https://www.youtube.com.eg",
  "http://www.youtube.com.eg",
  "youtube.com.kw",
  "www.youtube.com.kw",
  "https://www.youtube.com.kw",
  "http://www.youtube.com.kw",
  "youtube.com.qa",
  "www.youtube.com.qa",
  "https://www.youtube.com.qa",
  "http://www.youtube.com.qa",
  "youtube.com.om",
  "www.youtube.com.om",
  "https://www.youtube.com.om",
  "http://www.youtube.com.om",
  "youtube.com.bh",
  "www.youtube.com.bh",
  "https://www.youtube.com.bh",
  "http://www.youtube.com.bh"]

        if any(self.input_field.text.startswith(prefix) for prefix in prefixes):
                response = requests.post("http://176.56.54.46:1234/download_video", json={"url": self.input_field.text, "user": Config.get('postproc', 'email')})
                if response.status_code == 200:
                    data = response.json()
                    text = data.get("text")
                    Config.set('postproc', 'text', text)
                    Config.write()
                    self.input_field.text =''
                    sm.current = 'result'
                else:
                    print(f"Ошибка: {response.status_code}")
            
        else:
            Config.set('postproc', 'text', self.input_field.text)
            Config.write()
            self.input_field.text =''
            self.manager.current = 'result'




class ResultScreen(Screen):
    def __init__(self, **kwargs):
        super().__init__(**kwargs)
        self.name = 'result'
        

        with self.canvas:
            Color(0.9, 0.2, 0.2, 1)
            self.rect_top = Rectangle(pos=(0, 0), size=(self.width, self.height / 2))
            Color(0.9, 0.6, 0.2, 1)
            self.rect_bottom = Rectangle(pos=(0, self.height / 2), size=(self.width, self.height / 2))
            Color(0, 0, 0, 1)
            Line(points=[(self.width * 0.05, self.height * 0.05), (self.width * 0.1, self.height * 0.05),
                         (self.width * 0.1, self.height * 0.15), (self.width * 0.05, self.height * 0.15)])
            Line(points=[(self.width * 0.1, self.height * 0.05), (self.width * 0.1, self.height * 0.15),
                         (self.width * 0.15, self.height * 0.15), (self.width * 0.15, self.height * 0.05),
                         (self.width * 0.2, self.height * 0.05), (self.width * 0.2, self.height * 0.15)])
            Line(points=[(self.width * 0.2, self.height * 0.05), (self.width * 0.2, self.height * 0.15),
                         (self.width * 0.25, self.height * 0.15), (self.width * 0.25, self.height * 0.05)])

        title = Label(text='Education Assistant', font_size=24, bold=True,
                      pos_hint={'center_x': 0.5, 'center_y': 0.9}, color=[0, 0, 0, 1])
        self.add_widget(title)

        self.theme_label = Label(text='Тема: загрузка...', font_size='18sp', bold=True, color=(0.9, 0.2, 0.2, 1),
                                 size_hint=(None, None), size=(250, 35), pos_hint={'center_x': 0.5, 'top': 0.85})
        self.add_widget(self.theme_label)

        undo = Image(source='undo.png',
                      pos_hint={'center_x': 0.1, 'center_y': 0.9},
                      size_hint=(0.1, 0.1), allow_stretch=True)
        self.add_widget(undo)

        self.undo_button = Button(pos_hint={'center_x': 0.1, 'center_y': 0.9},
                           size_hint=(0.1, 0.1),
                           background_color=(1, 1, 1, 0),  # Прозрачный фон
                           border=(0, 0, 0, 0),  # Удаляем границы
                           color=(0, 0, 1, 1) # Цвет текста
                          )
        self.undo_button.bind(on_press=self.on_undo_button_pressed)
        self.add_widget(self.undo_button)

        folder = Image(source='folder.png',
                      pos_hint={'center_x': 0.9, 'center_y': 0.9},
                      size_hint=(0.1, 0.1), allow_stretch=True)
        self.add_widget(folder)

        

        self.folder_button = Button(pos_hint={'center_x': 0.9, 'center_y': 0.9},
                           size_hint=(0.1, 0.1),
                           background_color=(1, 1, 1, 0),  # Прозрачный фон
                           border=(0, 0, 0, 0),  # Удаляем границы
                           color=(0, 0, 1, 1) # Цвет текста
                          )
        self.folder_button.bind(on_press=self.on_folder_button_pressed)
        self.add_widget(self.folder_button)

        self.conspect_label = Label(text='Конспект:', font_size='18sp', bold=True, color=(0, 0, 0, 1),
                                 size_hint=(None, None), size=(250, 35), pos_hint={'center_x': 0.5, 'top': 0.8})
        self.add_widget(self.conspect_label)

        self.conspect = TextInput(text='Загрузка...', readonly = True, disabled = False, multiline = True,
                                      background_color=(1, 1, 1, 1), size_hint=(0.85, 0.28),
                                      pos_hint={'center_x': 0.5, 'top': 0.75})
        self.add_widget(self.conspect)

        self.plan_label = Label(text='План:', font_size='18sp', bold=True, color=(0, 0, 0, 1),
                                 size_hint=(None, None), size=(250, 35), pos_hint={'center_x': 0.5, 'top': 0.45})
        self.add_widget(self.plan_label)

        self.plan = TextInput(text='Загрузка...', readonly = True, disabled = False, multiline = True,
                                      background_color=(1, 1, 1, 1), size_hint=(0.85, 0.28),
                                      pos_hint={'center_x': 0.5, 'top': 0.4})
        self.add_widget(self.plan)

    
    def on_enter (self):
        if (Config.get('postproc', 'file') != ''):
            response = requests.post("http://176.56.54.46:1234/download", json={"user": Config.get('postproc', 'email'), "file_name": Config.get('postproc', 'file')})
            if response.status_code == 200:
                with open(f"{Config.get('postproc', 'file')}.json", "wb") as f:
                    data = response.json()
                    theme = data.get("theme")
                    sum_text = data.get("sum_text")
                    plan = data.get("plan")
                    self.theme_label.text = theme
                    self.conspect.text = sum_text
                    self.plan.text = plan
                    
                    Config.set('postproc', 'file', '')
                    Config.write()
            else:
                Config.set('postproc', 'file', '')
                Config.write()
                print(f"Ошибка скачивания файла: {response.status_code}")
        else:
            user = Config.get('postproc', 'email')
            text = Config.get('postproc', 'text')
            with open("output.txt", "w", encoding="utf-8") as file:
                file.write(text)
            with open("output.txt", "rb") as file:
                # Открыть файл в бинарном режиме
                response = requests.post("http://176.56.54.46:1234/upload_text", data={"user": user}, files={"file": file})
                if response.status_code == 200:
                    data = response.json()
                    theme = data.get("theme")
                    sum_text = data.get("sum_text")
                    plan = data.get("plan")
                    self.theme_label.text = theme
                    self.conspect.text = sum_text
                    self.plan.text = plan
                Config.set('postproc', 'text', '')
                Config.write()

        
    def on_undo_button_pressed(self, instance):
        self.theme_label.text = ''
        self.conspect.text = ''
        self.plan.text = ''
        self.manager.get_screen('main').delete_widget()
        self.manager.current = 'main'
        return

    def on_size(self, *args):
            self.rect_top.size = self.size
            self.rect_bottom.pos = (0, self.height / 2)
            self.rect_bottom.size = self.size

    def on_folder_button_pressed(self, instance):
        self.theme_label.text = ''
        self.conspect.text = ''
        self.plan.text = ''
        self.manager.get_screen('archive').delete_widget()

        self.manager.current = 'archive'
        return

        
        

class SettingsScreen(Screen):
    def __init__(self, **kwargs):
        super().__init__(**kwargs)
        self.name = 'settings'
        
        with self.canvas:
            Color(0.9, 0.2, 0.2, 1)
            self.rect_top = Rectangle(pos=(0, 0), size=(self.width, self.height / 2))
            Color(0.9, 0.6, 0.2, 1)
            self.rect_bottom = Rectangle(pos=(0, self.height / 2), size=(self.width, self.height / 2))
            Color(0, 0, 0, 1)
            Line(points=[(self.width * 0.05, self.height * 0.05), (self.width * 0.1, self.height * 0.05),
                         (self.width * 0.1, self.height * 0.15), (self.width * 0.05, self.height * 0.15)])
            Line(points=[(self.width * 0.1, self.height * 0.05), (self.width * 0.1, self.height * 0.15),
                         (self.width * 0.15, self.height * 0.15), (self.width * 0.15, self.height * 0.05),
                         (self.width * 0.2, self.height * 0.05), (self.width * 0.2, self.height * 0.15)])
            Line(points=[(self.width * 0.2, self.height * 0.05), (self.width * 0.2, self.height * 0.15),
                         (self.width * 0.25, self.height * 0.15), (self.width * 0.25, self.height * 0.05)])

        undo = Image(source='undo.png',
                      pos_hint={'center_x': 0.1, 'center_y': 0.9},
                      size_hint=(0.1, 0.1), allow_stretch=True)
        self.add_widget(undo)

        self.undo_button = Button(pos_hint={'center_x': 0.1, 'center_y': 0.9},
                           size_hint=(0.1, 0.1),
                           background_color=(1, 1, 1, 0),  # Прозрачный фон
                           border=(0, 0, 0, 0),  # Удаляем границы
                           color=(0, 0, 1, 1) # Цвет текста
                          )
        self.undo_button.bind(on_press=self.on_undo_button_pressed)
        self.add_widget(self.undo_button)

        # Title label
        self.title_label = Label(text='Education Assistant', font_size='24sp', bold=True, color=(0.9, 0.2, 0.2, 1),
                                 size_hint=(None, None), size=(250, 35), pos_hint={'center_x': 0.5, 'top': 0.95})
        self.add_widget(self.title_label)

        # Settings label
        self.settings_label = Label(text='Настройки', font_size='20sp', bold=True, color=(0.9, 0.2, 0.2, 1),
                                     size_hint=(None, None), size=(150, 30), pos_hint={'center_x': 0.5, 'top': 0.85})
        self.add_widget(self.settings_label)

        # Settings label
        self.settings_label = Label(text='Аккаунт', bold=True, color=(0.9, 0.2, 0.2, 1),
                                     size_hint=(None, None), size=(150, 30), pos_hint={'center_x': 0.5, 'top': 0.8})
        self.add_widget(self.settings_label)

        # Email 
        self.email_info = TextInput(text=Config.get('postproc', 'email'), multiline=False, readonly = True, disabled = True, padding=[10, 10, 10, 10],
                                      background_color=(1, 1, 1, 1), size_hint=(None, None), size=(250, 40),
                                      pos_hint={'center_x': 0.5, 'top': 0.75})
        self.add_widget(self.email_info)

        # Local storage label
        self.storage_label = Label(text='Хранилище локальных сохранений', bold=True, color=(0.9, 0.2, 0.2, 1),
                                    size_hint=(None, None), size=(250, 25), pos_hint={'center_x': 0.5, 'top': 0.6})
        self.add_widget(self.storage_label)

        # Path input
        self.path_input = TextInput(hint_text='В разработке...', multiline=False, disabled = True, padding=[10, 10, 10, 10],
                                      background_color=(1, 1, 1, 1), size_hint=(None, None), size=(250, 40),
                                      pos_hint={'center_x': 0.5, 'top': 0.55})
        self.path_input.bind(text = self.path_changed)
        self.add_widget(self.path_input)

        # Exit button
        self.exit_button = Button(text='Выйти', background_color=(1, 1, 0, 1), color=(0.9, 0.2, 0.2, 1), font_size='18sp',
                                   size_hint=(None, None), size=(150, 50), pos_hint={'center_x': 0.5, 'top': 0.35})
        self.exit_button.bind(on_press=self.on_exit_button_pressed)
        self.add_widget(self.exit_button)

    def on_enter(self):
        try:
                self.path_input.text = Config.get('postproc', 'path')
                self.email_info.text=Config.get('postproc', 'email')
        except:
            pass

    def on_size(self, *args):
        self.rect_top.size = self.size
        self.rect_bottom.pos = (0, self.height / 2)
        self.rect_bottom.size = self.size

    def path_changed(self, instance, text):
        Config.set('postproc', 'path', self.path_input.text)
        Config.write()

    def on_undo_button_pressed(self, instance):
        self.manager.get_screen('main').delete_widget()
        self.manager.current = 'main'
        return

    def on_exit_button_pressed(self, instance):
        Config.set('postproc', 'email', '')
        Config.write()
        Config.set('postproc', 'password', '')
        Config.write()
        self.manager.current = 'reg'




class ArchiveScreen(Screen):
    
    def __init__(self, **kwargs):
                super(ArchiveScreen, self).__init__(**kwargs)

                self.name = 'archive'

                with self.canvas:
                    Color(0.9, 0.2, 0.2, 1)
                    self.rect_top = Rectangle(pos=(0, 0), size=(self.width, self.height / 2))
                    Color(0.9, 0.6, 0.2, 1)
                    self.rect_bottom = Rectangle(pos=(0, self.height / 2), size=(self.width, self.height / 2))
                    Color(0, 0, 0, 1)
                    Line(points=[(self.width * 0.05, self.height * 0.05), (self.width * 0.1, self.height * 0.05),
                         (self.width * 0.1, self.height * 0.15), (self.width * 0.05, self.height * 0.15)])
                    Line(points=[(self.width * 0.1, self.height * 0.05), (self.width * 0.1, self.height * 0.15),
                         (self.width * 0.15, self.height * 0.15), (self.width * 0.15, self.height * 0.05),
                         (self.width * 0.2, self.height * 0.05), (self.width * 0.2, self.height * 0.15)])
                    Line(points=[(self.width * 0.2, self.height * 0.05), (self.width * 0.2, self.height * 0.15),
                         (self.width * 0.25, self.height * 0.15), (self.width * 0.25, self.height * 0.05)])


                self.title_label = Label(text='Education Assistant', font_size='24sp', bold=True, color=(0.9, 0.2, 0.2, 1),
                                 size_hint=(None, None), size=(250, 35), pos_hint={'center_x': 0.5, 'top': 0.95})
                self.add_widget(self.title_label)

                self.rv = RV()
                self.add_widget(self.rv)
                

                self.undo = Image(source='undo.png',
                      pos_hint={'center_x': 0.1, 'center_y': 0.9},
                      size_hint=(0.1, 0.1), allow_stretch=True)
                self.add_widget(self.undo)

                self.undo_button = Button(pos_hint={'center_x': 0.1, 'center_y': 0.9},
                           size_hint=(0.1, 0.1),
                           background_color=(1, 1, 1, 0),  # Прозрачный фон
                           border=(0, 0, 0, 0),  # Удаляем границы
                           color=(0, 0, 1, 1) # Цвет текста
                          )
                self.undo_button.bind(on_press=self.on_undo_button_pressed)
                self.add_widget(self.undo_button)

                
    
    def delete_widget(self):
            self.remove_widget(self.undo)
            self.remove_widget(self.undo_button)
            self.remove_widget(self.rv)
            self.rv = RV()
            self.add_widget(self.rv)
            self.add_widget(self.undo)
            self.add_widget(self.undo_button)
        
        

    def on_size(self, *args):
        self.rect_top.size = self.size
        self.rect_bottom.pos = (0, self.height / 2)
        self.rect_bottom.size = self.size

    def on_undo_button_pressed(self, instance):
        self.manager.get_screen('main').delete_widget()
        self.manager.current = 'main'

        


class App(App):
    def build(self):

        def func():
            f = open('result.ini', 'w') # и в функции записывать в файл значение
            b = f.write("The function was completed")
            f.close()
            Config.set('postproc', 'email', '')
            Config.write()
            Config.set('postproc', 'password', '')
            Config.write()

        f = open('result.ini', 'r') #при каждом запуске читаются данные
        a = f.read()
        f.close()

        if a == "": #если файл пустой - выполнить функцию
            func()

        sc1 = RegScreen(name='reg')
        sc2 = AuthScreen(name='auth')
        sc3 = MainScreen(name='main')
        sc4 = SettingsScreen(name='settings')
        sc5 = ArchiveScreen(name='archive')
        sc6 = ResultScreen(name='result')
        
        

        if Config.get('postproc', 'email') != '' and Config.get('postproc', 'password') != '':
            user = User(email=Config.get('postproc', 'email'), password=Config.get('postproc', 'password'))

            response = requests.post("http://176.56.54.46:1234/auth", json=user.dict())

            if response.status_code == 200:
                sm.add_widget(sc3)
                sm.add_widget(sc4)
                sm.add_widget(sc2)
                sm.add_widget(sc1)
                sm.add_widget(sc5)
                sm.add_widget(sc6)
                

            else:
                sm.add_widget(sc1)
                sm.add_widget(sc4)
                sm.add_widget(sc3)
                sm.add_widget(sc2)
                sm.add_widget(sc5)
                sm.add_widget(sc6)

        else:
            sm.add_widget(sc1)
            sm.add_widget(sc4)
            sm.add_widget(sc3)
            sm.add_widget(sc2)
            sm.add_widget(sc5)
            sm.add_widget(sc6)

        return sm

def split_recog(audio_file):
        
        max_duration=61
        audio = AudioSegment.from_file(audio_file)
        text = ""
        r = sr.Recognizer()
        result = ""
        fragments = []
        start_time = 0
        end_time = 0

        while end_time < len(audio):
            end_time = min(start_time + max_duration * 1000, len(audio))
            audio_fragment = audio[start_time:end_time]

    # Используем распознавание речи для определения конца слова
            with sr.AudioFile(audio_fragment.export("temp.wav", format="wav")) as source:
                audio_data = r.record(source)
                try:
                    text = r.recognize_google(audio_data, language="ru-RU")
        # Находим индекс последнего пробела в тексте, чтобы определить конец слова
                    last_space_index = text.rfind("  ")
        # Если пробел найден, сдвигаем время окончания фрагмента
                    if last_space_index != -1:
          # Исправленный расчет времени окончания:
                        end_time = start_time + (last_space_index + 1) * 1000 / len(text) * len(audio_fragment)

          # Проверка, чтобы не выйти за пределы аудиофайла
                        end_time = min(end_time, len(audio))

                except sr.UnknownValueError:
                    pass
                except sr.RequestError as e:
                    print(f"Ошибка распознавания речи: {e}")

            result = result + text + " "
    # Сохраняем фрагмент
            fragment_name = f"{audio_file[:-4]}_{start_time // 1000}_{end_time // 1000}.wav"
            audio_fragment.export(fragment_name, format="wav")
            fragments.append(fragment_name)

    # Переходим к следующему фрагменту, **только если end_time не достиг конца аудиофайла**
            if end_time < len(audio):
                start_time = end_time
        print(result)
        Config.set('postproc', 'text', result)
        Config.write()
        Clock.schedule_once(update_screen, 0)

def extract_audio_from_video(video_file, audio_file):
    
    video = VideoFileClip(video_file)
    audio = video.audio
    audio.write_audiofile(audio_file)
    split_recog(audio_file)
    
def update_screen(self):
    sm.current = 'result'

class User(BaseModel):
    email: str
    password: str


if __name__ == '__main__':
    App().run()
